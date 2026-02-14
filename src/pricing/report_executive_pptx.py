from __future__ import annotations

"""
Generate executive Markdown and PPTX deliverables from a pricing config.
"""

from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
import copy
import json
import os
from typing import Any, Mapping

import pandas as pd
import yaml

from .config import read_loading_parameters
from .diagnostics import build_execution_context, build_run_summary
from .endowment import LoadingFunctionParams, calc_loading_parameters
from .paths import resolve_base_dir_from_config
from .profit_test import DEFAULT_LAPSE_RATE, DEFAULT_VALUATION_INTEREST, run_profit_test
from .report_feasibility import build_feasibility_report


EXPENSE_SCALE_COLUMNS = (
    "acq_var_total",
    "acq_fixed_total",
    "maint_var_total",
    "maint_fixed_total",
    "coll_var_total",
    "overhead_total",
)


@dataclass(frozen=True)
class ExecutiveReportOutputs:
    pptx_path: Path
    markdown_path: Path
    run_summary_path: Path
    feasibility_deck_path: Path
    cashflow_chart_path: Path
    premium_chart_path: Path


def _resolve_output_path(base_dir: Path, path: Path | None, default: str) -> Path:
    target = Path(default) if path is None else Path(path)
    return target if target.is_absolute() else (base_dir / target)


def _require_matplotlib():
    try:
        import matplotlib.pyplot as plt
    except ModuleNotFoundError as exc:  # pragma: no cover - depends on runtime env
        raise RuntimeError(
            "matplotlib is required for report-executive-pptx. "
            "Install with: python -m pip install matplotlib"
        ) from exc
    return plt


def _configure_plot_font_for_language(plt, language: str) -> None:
    if language != "ja":
        return
    try:
        from matplotlib import font_manager
    except Exception:  # pragma: no cover - defensive
        return

    available = {font.name for font in font_manager.fontManager.ttflist}
    candidates = [
        "Yu Gothic",
        "Meiryo",
        "MS Gothic",
        "Noto Sans CJK JP",
        "IPAexGothic",
    ]
    for name in candidates:
        if name in available:
            plt.rcParams["font.family"] = [name]
            plt.rcParams["axes.unicode_minus"] = False
            return


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:  # pragma: no cover - depends on runtime env
        raise RuntimeError(
            "python-pptx is required for report-executive-pptx. "
            "Install with: python -m pip install python-pptx"
        ) from exc
    return Presentation, RGBColor, Inches, Pt


def _fmt_pct(value: float, digits: int = 2) -> str:
    return f"{value * 100:.{digits}f}%"


def _fmt_jpy(value: float) -> str:
    return f"JPY {value:,.0f}"


def _validate_language(language: str) -> str:
    lang = str(language).strip().lower()
    if lang not in {"ja", "en"}:
        raise ValueError(f"Unsupported language: {language}. Use 'ja' or 'en'.")
    return lang


def _scenario_label(name: str, language: str) -> str:
    if language == "en":
        return name
    mapping = {
        "base": "ベース",
        "interest_down_10pct": "金利-10%",
        "interest_up_10pct": "金利+10%",
        "lapse_down_10pct": "解約率-10%",
        "lapse_up_10pct": "解約率+10%",
        "expense_down_10pct": "事業費-10%",
        "expense_up_10pct": "事業費+10%",
    }
    return mapping.get(name, name)


def _constraint_label(name: str, language: str) -> str:
    if language == "en":
        return name
    mapping = {
        "irr_hard": "IRR下限",
        "nbv_hard": "NBV下限",
        "loading_surplus_hard": "負荷余剰下限",
        "loading_surplus_ratio_hard": "負荷余剰率下限",
        "premium_to_maturity_hard_max": "PTM上限",
    }
    return mapping.get(name, name)


def _aggregate_cashflow(batch_result) -> pd.DataFrame:
    if not batch_result.results:
        raise ValueError("No model point results available.")
    frames = [res.cashflow for res in batch_result.results]
    all_cashflow = pd.concat(frames, ignore_index=True)
    agg = (
        all_cashflow.groupby("t", as_index=False)[
            [
                "premium_income",
                "investment_income",
                "death_benefit",
                "surrender_benefit",
                "expenses_total",
                "reserve_change",
                "net_cf",
            ]
        ]
        .sum()
        .sort_values("t")
    )
    agg["year"] = agg["t"].astype(int) + 1
    agg["benefit_outgo"] = -(agg["death_benefit"] + agg["surrender_benefit"])
    agg["expense_outgo"] = -agg["expenses_total"]
    agg["reserve_change_outgo"] = -agg["reserve_change"]
    return agg


def _plot_cashflow_by_profit_source(
    agg: pd.DataFrame,
    out_path: Path,
    *,
    language: str,
) -> Path:
    plt = _require_matplotlib()
    _configure_plot_font_for_language(plt, language)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    years = agg["year"].astype(int).tolist()
    if language == "ja":
        positive = [
            ("premium_income", "保険料収入", "#0b5fa5"),
            ("investment_income", "運用収益", "#2a9d8f"),
        ]
        negative = [
            ("benefit_outgo", "保険金・解約返戻金", "#d1495b"),
            ("expense_outgo", "事業費", "#f4a261"),
            ("reserve_change_outgo", "責任準備金増減", "#6c757d"),
        ]
        net_cf_label = "純キャッシュフロー"
        chart_title = "年度別キャッシュフロー（利源別、全モデルポイント合算）"
        x_label = "保険年度"
    else:
        positive = [
            ("premium_income", "Premium Income", "#0b5fa5"),
            ("investment_income", "Investment Income", "#2a9d8f"),
        ]
        negative = [
            ("benefit_outgo", "Benefit Outgo", "#d1495b"),
            ("expense_outgo", "Expense Outgo", "#f4a261"),
            ("reserve_change_outgo", "Reserve Change", "#6c757d"),
        ]
        net_cf_label = "Net CF"
        chart_title = "Yearly Cashflow by Profit Source (All Model Points)"
        x_label = "Policy Year"

    fig, ax = plt.subplots(figsize=(12, 5), dpi=150)
    pos_base = [0.0 for _ in years]
    for col, label, color in positive:
        vals = [max(float(v), 0.0) for v in agg[col].tolist()]
        ax.bar(years, vals, bottom=pos_base, label=label, color=color, width=0.8)
        pos_base = [b + v for b, v in zip(pos_base, vals)]

    neg_base = [0.0 for _ in years]
    for col, label, color in negative:
        vals = [min(float(v), 0.0) for v in agg[col].tolist()]
        ax.bar(years, vals, bottom=neg_base, label=label, color=color, width=0.8)
        neg_base = [b + v for b, v in zip(neg_base, vals)]

    ax.plot(years, agg["net_cf"], color="#111111", linewidth=2.0, label=net_cf_label)
    ax.set_title(chart_title)
    ax.set_xlabel(x_label)
    ax.set_ylabel("JPY")
    if years:
        tick_step = max(1, len(years) // 12)
        ax.set_xticks(years[::tick_step])
    ax.axhline(0.0, color="#333333", linewidth=0.8)
    ax.grid(axis="y", alpha=0.25)
    ax.legend(ncol=3, fontsize=8, frameon=False, loc="upper right")
    fig.tight_layout()
    fig.savefig(out_path)
    plt.close(fig)
    return out_path


def _plot_annual_premium_by_model_point(
    summary_df: pd.DataFrame,
    out_path: Path,
    *,
    language: str,
) -> Path:
    plt = _require_matplotlib()
    _configure_plot_font_for_language(plt, language)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    chart_df = summary_df.sort_values("gross_annual_premium", ascending=True)
    fig, ax = plt.subplots(figsize=(10, 4.5), dpi=150)
    ax.barh(
        chart_df["model_point"].tolist(),
        chart_df["gross_annual_premium"].astype(float).tolist(),
        color="#0b5fa5",
    )
    if language == "ja":
        ax.set_title("モデルポイント別 年払保険料P")
        ax.set_xlabel("年払営業保険料（JPY）")
    else:
        ax.set_title("Annual Premium P by Model Point")
        ax.set_xlabel("Gross Annual Premium (JPY)")
    ax.grid(axis="x", alpha=0.25)
    fig.tight_layout()
    fig.savefig(out_path)
    plt.close(fig)
    return out_path


def _constraint_status_rows(run_summary: Mapping[str, Any]) -> list[dict[str, object]]:
    status_by_type: dict[str, dict[str, object]] = {}
    for model_point in run_summary["model_points"]:
        for entry in model_point["constraints"]:
            key = str(entry["type"])
            current = status_by_type.get(key)
            if current is None:
                status_by_type[key] = {
                    "constraint": key,
                    "threshold": float(entry["threshold"]),
                    "min_gap": float(entry["gap"]),
                    "worst_model_point": str(model_point["model_point"]),
                    "all_ok": bool(entry["ok"]),
                }
                continue
            gap = float(entry["gap"])
            if gap < float(current["min_gap"]):
                current["min_gap"] = gap
                current["worst_model_point"] = str(model_point["model_point"])
            if not bool(entry["ok"]):
                current["all_ok"] = False

    rows = list(status_by_type.values())
    rows.sort(key=lambda row: str(row["constraint"]))
    return rows


def _resolve_company_expense_path(config: Mapping[str, object], base_dir: Path) -> Path | None:
    profit_test_cfg = config.get("profit_test", {})
    if not isinstance(profit_test_cfg, Mapping):
        return None
    expense_cfg = profit_test_cfg.get("expense_model", {})
    if not isinstance(expense_cfg, Mapping):
        return None
    raw = expense_cfg.get("company_data_path")
    if not isinstance(raw, str):
        return None
    path = Path(raw)
    return path if path.is_absolute() else (base_dir / path)


def _scale_company_expense_file(original_path: Path, factor: float, scaled_path: Path) -> Path:
    df = pd.read_csv(original_path)
    for col in EXPENSE_SCALE_COLUMNS:
        if col in df.columns:
            scaled = df[col].astype(float) * float(factor)
            if (scaled < 0.0).any():
                raise ValueError(f"Negative planned expense assumptions are not allowed: {col}")
            df[col] = scaled
    scaled_path.parent.mkdir(parents=True, exist_ok=True)
    df.to_csv(scaled_path, index=False)
    return scaled_path


def _scenario_summary(name: str, config: dict, base_dir: Path) -> dict[str, object]:
    result = run_profit_test(config, base_dir=base_dir)
    summary = build_run_summary(config, result, source=f"sensitivity:{name}")
    metrics = summary["summary"]
    return {
        "scenario": name,
        "min_irr": float(metrics["min_irr"]),
        "min_nbv": float(metrics["min_nbv"]),
        "min_loading_surplus_ratio": float(metrics["min_loading_surplus_ratio"]),
        "max_premium_to_maturity": float(metrics["max_premium_to_maturity"]),
        "violation_count": int(metrics["violation_count"]),
    }


def _build_sensitivity_rows(config: dict, base_dir: Path, temp_dir: Path) -> list[dict[str, object]]:
    scenarios: list[dict[str, object]] = []
    scenarios.append(_scenario_summary("base", config, base_dir))

    pricing_cfg = config.get("pricing", {})
    if isinstance(pricing_cfg, Mapping):
        interest_cfg = pricing_cfg.get("interest", {})
        if isinstance(interest_cfg, Mapping):
            flat_rate = float(interest_cfg.get("flat_rate", 0.0))
            for factor, label in ((0.9, "interest_down_10pct"), (1.1, "interest_up_10pct")):
                scenario_cfg = copy.deepcopy(config)
                scenario_cfg["pricing"]["interest"]["flat_rate"] = flat_rate * factor
                profit_test_cfg = scenario_cfg.setdefault("profit_test", {})
                valuation = float(
                    profit_test_cfg.get("valuation_interest_rate", DEFAULT_VALUATION_INTEREST)
                )
                profit_test_cfg["valuation_interest_rate"] = valuation * factor
                scenarios.append(_scenario_summary(label, scenario_cfg, base_dir))

    lapse_base = float(
        config.get("profit_test", {}).get("lapse_rate", DEFAULT_LAPSE_RATE)
        if isinstance(config.get("profit_test", {}), Mapping)
        else DEFAULT_LAPSE_RATE
    )
    for factor, label in ((0.9, "lapse_down_10pct"), (1.1, "lapse_up_10pct")):
        scenario_cfg = copy.deepcopy(config)
        scenario_cfg.setdefault("profit_test", {})["lapse_rate"] = lapse_base * factor
        scenarios.append(_scenario_summary(label, scenario_cfg, base_dir))

    expense_path = _resolve_company_expense_path(config, base_dir)
    if expense_path is not None and expense_path.is_file():
        for factor, label in ((0.9, "expense_down_10pct"), (1.1, "expense_up_10pct")):
            scenario_cfg = copy.deepcopy(config)
            scaled = _scale_company_expense_file(
                expense_path,
                factor,
                temp_dir / f"{expense_path.stem}_{label}.csv",
            )
            scenario_cfg.setdefault("profit_test", {}).setdefault("expense_model", {})[
                "company_data_path"
            ] = str(scaled.resolve())
            scenarios.append(_scenario_summary(label, scenario_cfg, base_dir))
    return scenarios


def _get_loading_params(config: Mapping[str, object]) -> LoadingFunctionParams | None:
    return read_loading_parameters(config)


def _loading_calculation_rows(batch_result, params: LoadingFunctionParams | None) -> list[dict[str, object]]:
    rows: list[dict[str, object]] = []
    for result in batch_result.results:
        point = result.model_point
        age_delta = float(point.issue_age - 30)
        term_delta = float(point.term_years - 10)
        sex_indicator = 1.0 if point.sex == "female" else 0.0
        if params is None:
            alpha = float(result.loadings.alpha)
            beta = float(result.loadings.beta)
            gamma = float(result.loadings.gamma)
            alpha_expr = f"{alpha:.6f} (fixed)"
            beta_expr = f"{beta:.6f} (fixed)"
            gamma_expr = f"{gamma:.6f} (fixed)"
        else:
            generated = calc_loading_parameters(
                params=params,
                issue_age=point.issue_age,
                term_years=point.term_years,
                sex=point.sex,
            )
            alpha = float(generated.alpha)
            beta = float(generated.beta)
            gamma_raw = params.g0 + params.g_term * term_delta
            gamma = float(generated.gamma)
            alpha_expr = (
                f"{alpha:.6f} = {params.a0:.6f} + ({params.a_age:.6f}*{age_delta:.1f}) + "
                f"({params.a_term:.6f}*{term_delta:.1f}) + ({params.a_sex:.6f}*{sex_indicator:.1f})"
            )
            beta_expr = (
                f"{beta:.6f} = {params.b0:.6f} + ({params.b_age:.6f}*{age_delta:.1f}) + "
                f"({params.b_term:.6f}*{term_delta:.1f}) + ({params.b_sex:.6f}*{sex_indicator:.1f})"
            )
            gamma_expr = (
                f"{gamma:.6f} = clamp({params.g0:.6f} + ({params.g_term:.6f}*{term_delta:.1f})"
                f" = {gamma_raw:.6f}, 0.0, 0.5)"
            )
        rows.append(
            {
                "model_point": point.model_point_id,
                "sex": point.sex,
                "issue_age": point.issue_age,
                "term_years": point.term_years,
                "age_delta": age_delta,
                "term_delta": term_delta,
                "sex_indicator": sex_indicator,
                "alpha": alpha,
                "beta": beta,
                "gamma": gamma,
                "alpha_expr": alpha_expr,
                "beta_expr": beta_expr,
                "gamma_expr": gamma_expr,
            }
        )
    return rows


def _build_markdown_report(
    *,
    config_path: Path,
    markdown_path: Path,
    run_summary: Mapping[str, Any],
    batch_result,
    feasibility_deck: Mapping[str, Any],
    sensitivity_rows: list[dict[str, object]],
    cashflow_chart_path: Path,
    premium_chart_path: Path,
    language: str,
) -> str:
    language = _validate_language(language)
    is_ja = language == "ja"
    summary = run_summary["summary"]
    cashflow_rel = Path(os.path.relpath(cashflow_chart_path, markdown_path.parent)).as_posix()
    premium_rel = Path(os.path.relpath(premium_chart_path, markdown_path.parent)).as_posix()
    params = _get_loading_params(
        yaml.safe_load(config_path.read_text(encoding="utf-8"))
    )
    loading_rows = _loading_calculation_rows(batch_result, params)
    summary_df = batch_result.summary.sort_values("model_point")
    constraint_rows = _constraint_status_rows(run_summary)

    lines: list[str] = []
    report_title = "実現可能性レポート" if is_ja else "Feasibility Report"
    generated_label = "作成日" if is_ja else "generated"
    lines.append(f"# {report_title} ({config_path.name}, {generated_label} {datetime.now(timezone.utc).date().isoformat()})")
    lines.append("")
    lines.append("## サマリー" if is_ja else "## Summary")
    if is_ja:
        lines.append(
            f"- 制約評価: `violation_count={summary['violation_count']}` / `{summary['model_point_count']}`モデルポイント。"
        )
        lines.append(
            f"- 収益性下限: `min_irr={summary['min_irr']:.6f}`（現行設定での最悪点）。"
        )
        lines.append(f"- NBV下限: `min_nbv={summary['min_nbv']:.2f}` JPY。")
        lines.append(
            f"- 十分性下限: `min_loading_surplus_ratio={summary['min_loading_surplus_ratio']:.6f}`。"
        )
        lines.append(
            f"- 健全性上限: `max_premium_to_maturity={summary['max_premium_to_maturity']:.6f}`。"
        )
    else:
        lines.append(
            f"- Constraint status: `violation_count={summary['violation_count']}` across `{summary['model_point_count']}` model points."
        )
        lines.append(
            f"- Profitability floor: `min_irr={summary['min_irr']:.6f}` (worst model point at current setting)."
        )
        lines.append(f"- NBV floor: `min_nbv={summary['min_nbv']:.2f}` JPY.")
        lines.append(
            f"- Adequacy floor: `min_loading_surplus_ratio={summary['min_loading_surplus_ratio']:.6f}`."
        )
        lines.append(
            f"- Soundness ceiling: `max_premium_to_maturity={summary['max_premium_to_maturity']:.6f}`."
        )
    lines.append("")
    lines.append("## 価格提案（モデルポイント別P）" if is_ja else "## Pricing Recommendation (P by Model Point)")
    if is_ja:
        lines.append("|model_point|年払保険料P|月払保険料|irr|nbv|premium_to_maturity|")
    else:
        lines.append("|model_point|gross_annual_premium|monthly_premium|irr|nbv|premium_to_maturity|")
    lines.append("|---|---:|---:|---:|---:|---:|")
    for row in summary_df.itertuples(index=False):
        lines.append(
            "|"
            f"{row.model_point}|{int(row.gross_annual_premium)}|{int(row.monthly_premium)}|"
            f"{row.irr:.6f}|{row.new_business_value:.2f}|{row.premium_to_maturity_ratio:.6f}|"
        )
    lines.append("")
    lines.append("## 制約ステータス" if is_ja else "## Constraint Status")
    if is_ja:
        lines.append("|constraint|min_gap|worst_model_point|status|")
    else:
        lines.append("|constraint|min_gap|worst_model_point|status|")
    lines.append("|---|---:|---|---|")
    for row in constraint_rows:
        status = "適合" if (is_ja and bool(row["all_ok"])) else "要対応" if is_ja else "PASS" if bool(row["all_ok"]) else "FAIL"
        constraint_name = _constraint_label(str(row["constraint"]), language)
        lines.append(
            f"|{constraint_name}|{float(row['min_gap']):.6f}|{row['worst_model_point']}|{status}|"
        )
    lines.append("")
    lines.append("## ローディング式と係数" if is_ja else "## Loading Formula and Coefficients")
    lines.append("- `gross_rate = (net_rate + alpha / a + beta) / (1 - gamma)`")
    if params is None:
        lines.append("- ローディングモード: 固定 `loading_alpha_beta_gamma`。" if is_ja else "- Loading mode: fixed `loading_alpha_beta_gamma`.")
    else:
        lines.append("- ローディングモード: `loading_parameters`。" if is_ja else "- Loading mode: `loading_parameters`.")
        lines.append("|係数|値|" if is_ja else "|coefficient|value|")
        lines.append("|---|---:|")
        for key in ("a0", "a_age", "a_term", "a_sex", "b0", "b_age", "b_term", "b_sex", "g0", "g_term"):
            lines.append(f"|{key}|{getattr(params, key):.6f}|")
    lines.append("")
    lines.append(
        "## モデルポイント別 alpha/beta/gamma 計算（中間計算付き）"
        if is_ja
        else "## Per-model-point alpha/beta/gamma Calculations (with intermediate steps)"
    )
    for row in loading_rows:
        if is_ja:
            lines.append(
                f"### {row['model_point']} (age={row['issue_age']}, term={row['term_years']}, sex={row['sex']})"
            )
            lines.append(
                f"- 差分: `age_delta={row['age_delta']:.1f}`, `term_delta={row['term_delta']:.1f}`, `sex_indicator={row['sex_indicator']:.1f}`"
            )
            lines.append(f"- alpha: `{row['alpha_expr']}`")
            lines.append(f"- beta: `{row['beta_expr']}`")
            lines.append(f"- gamma: `{row['gamma_expr']}`")
        else:
            lines.append(
                f"### {row['model_point']} (age={row['issue_age']}, term={row['term_years']}, sex={row['sex']})"
            )
            lines.append(
                f"- deltas: `age_delta={row['age_delta']:.1f}`, `term_delta={row['term_delta']:.1f}`, `sex_indicator={row['sex_indicator']:.1f}`"
            )
            lines.append(f"- alpha: `{row['alpha_expr']}`")
            lines.append(f"- beta: `{row['beta_expr']}`")
            lines.append(f"- gamma: `{row['gamma_expr']}`")
    lines.append("")
    lines.append("## 年度別キャッシュフロー（利源別）" if is_ja else "## Yearly Cashflow by Profit Source")
    lines.append(f"![{'年度別キャッシュフロー（利源別）' if is_ja else 'Yearly Cashflow by Profit Source'}]({cashflow_rel})")
    lines.append("")
    lines.append(f"![{'モデルポイント別 年払保険料' if is_ja else 'Annual Premium by Model Point'}](" + premium_rel + ")")
    lines.append("")
    lines.append("## 感応度サマリー" if is_ja else "## Sensitivity Summary")
    if is_ja:
        lines.append("|scenario|min_irr|min_nbv|min_loading_surplus_ratio|max_premium_to_maturity|violation_count|")
    else:
        lines.append("|scenario|min_irr|min_nbv|min_loading_surplus_ratio|max_premium_to_maturity|violation_count|")
    lines.append("|---|---:|---:|---:|---:|---:|")
    for row in sensitivity_rows:
        scenario_name = _scenario_label(str(row["scenario"]), language)
        lines.append(
            f"|{scenario_name}|{float(row['min_irr']):.6f}|{float(row['min_nbv']):.2f}|"
            f"{float(row['min_loading_surplus_ratio']):.6f}|{float(row['max_premium_to_maturity']):.6f}|"
            f"{int(row['violation_count'])}|"
        )
    lines.append("")
    lines.append("## Feasibility Deck メタ情報" if is_ja else "## Feasibility Deck Meta")
    scan = feasibility_deck["meta"]["scan"]
    if is_ja:
        lines.append(
            f"- 掃引範囲: `r_start={scan['r_start']}`, `r_end={scan['r_end']}`, `r_step={scan['r_step']}`, `irr_threshold={scan['irr_threshold']}`"
        )
    else:
        lines.append(
            f"- sweep range: `r_start={scan['r_start']}`, `r_end={scan['r_end']}`, `r_step={scan['r_step']}`, `irr_threshold={scan['irr_threshold']}`"
        )
    lines.append("")
    lines.append("## 再現手順" if is_ja else "## Reproducibility")
    lines.append("```powershell")
    lines.append("python -m pytest -q")
    lines.append(f"python -m pricing.cli run {config_path.as_posix()}")
    lines.append(
        "python -m pricing.cli report-feasibility "
        f"{config_path.as_posix()} --r-start 1.00 --r-end 1.08 --r-step 0.005 --irr-threshold 0.02 "
        "--out out/feasibility_deck_executive.yaml"
    )
    lines.append(
        "python -m pricing.cli report-executive-pptx "
        f"{config_path.as_posix()} --out reports/executive_pricing_deck.pptx "
        f"--md-out reports/feasibility_report.md --lang {language}"
    )
    lines.append("```")
    return "\n".join(lines)


def _set_text_style(text_frame, *, size, bold=False, color=None):
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = size
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _write_executive_pptx(
    *,
    out_path: Path,
    run_summary: Mapping[str, Any],
    summary_df: pd.DataFrame,
    constraint_rows: list[dict[str, object]],
    sensitivity_rows: list[dict[str, object]],
    cashflow_chart_path: Path,
    premium_chart_path: Path,
    config_path: Path,
    language: str,
) -> Path:
    language = _validate_language(language)
    is_ja = language == "ja"
    Presentation, RGBColor, Inches, Pt = _require_pptx()
    prs = Presentation()

    def add_title(slide, title: str, subtitle: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = title
        _set_text_style(title_tf, size=Pt(30), bold=True, color=RGBColor(11, 95, 165))
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.4))
            sub_tf = sub_box.text_frame
            sub_tf.text = subtitle
            _set_text_style(sub_tf, size=Pt(13), bold=False, color=RGBColor(80, 80, 80))

    summary = run_summary["summary"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    summary_title = "エグゼクティブサマリー" if is_ja else "Executive Summary"
    summary_subtitle = (
        f"設定: {config_path.name} | 作成日: {datetime.now(timezone.utc).date().isoformat()}"
        if is_ja
        else f"Config: {config_path.name} | Generated: {datetime.now(timezone.utc).date().isoformat()}"
    )
    add_title(
        slide,
        summary_title,
        summary_subtitle,
    )
    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.8), Inches(4.8))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    bullets = (
        [
            "提案: 2枚目の価格テーブルP（年払保険料）を承認してください。",
            f"制約結果: violation_count={summary['violation_count']} / model_points={summary['model_point_count']}。",
            f"最小IRR={_fmt_pct(float(summary['min_irr']), 2)}、最小NBV={_fmt_jpy(float(summary['min_nbv']))}。",
            f"最大premium-to-maturity={float(summary['max_premium_to_maturity']):.6f}。",
            "全ての数値主張は out/run_summary_executive.json と out/feasibility_deck_executive.yaml で再現可能です。",
        ]
        if is_ja
        else [
            "Recommendation: approve pricing table P (annual premium) shown in Slide 2.",
            f"Constraint result: violation_count={summary['violation_count']} / model_points={summary['model_point_count']}.",
            f"Minimum IRR={_fmt_pct(float(summary['min_irr']), 2)}, minimum NBV={_fmt_jpy(float(summary['min_nbv']))}.",
            f"Maximum premium-to-maturity={float(summary['max_premium_to_maturity']):.6f}.",
            "All quantitative claims are reproducible from out/run_summary_executive.json and out/feasibility_deck_executive.yaml.",
        ]
    )
    tf.text = bullets[0]
    _set_text_style(tf, size=Pt(20), color=RGBColor(20, 20, 20))
    for item in bullets[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.runs[0].font.size = Pt(19)
        p.runs[0].font.color.rgb = RGBColor(20, 20, 20)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "価格提案" if is_ja else "Pricing Recommendation",
        "モデルポイント別 年払保険料P" if is_ja else "Annual Premium P by Model Point",
    )
    table_rows = len(summary_df) + 1
    table = slide.shapes.add_table(table_rows, 6, Inches(0.4), Inches(1.5), Inches(8.1), Inches(4.7)).table
    headers = (
        ["モデルポイント", "年払P", "月払P", "IRR", "NBV", "PTM"]
        if is_ja
        else ["Model Point", "Annual P", "Monthly P", "IRR", "NBV", "PTM"]
    )
    for col, name in enumerate(headers):
        table.cell(0, col).text = name
    for idx, row in enumerate(summary_df.itertuples(index=False), start=1):
        table.cell(idx, 0).text = str(row.model_point)
        table.cell(idx, 1).text = f"{int(row.gross_annual_premium):,}"
        table.cell(idx, 2).text = f"{int(row.monthly_premium):,}"
        table.cell(idx, 3).text = _fmt_pct(float(row.irr), 2)
        table.cell(idx, 4).text = f"{float(row.new_business_value):,.0f}"
        table.cell(idx, 5).text = f"{float(row.premium_to_maturity_ratio):.4f}"
    slide.shapes.add_picture(str(premium_chart_path), Inches(8.8), Inches(1.6), width=Inches(4.3))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "制約ステータス" if is_ja else "Constraint Status",
        "十分性 / 収益性 / 健全性" if is_ja else "Adequacy / Profitability / Soundness",
    )
    c_table = slide.shapes.add_table(
        len(constraint_rows) + 1, 4, Inches(0.4), Inches(1.5), Inches(12.0), Inches(4.0)
    ).table
    c_headers = (
        ["制約", "最小ギャップ", "最悪モデルポイント", "判定"]
        if is_ja
        else ["Constraint", "Min Gap", "Worst Model Point", "Status"]
    )
    for col, name in enumerate(c_headers):
        c_table.cell(0, col).text = name
    for idx, row in enumerate(constraint_rows, start=1):
        c_table.cell(idx, 0).text = _constraint_label(str(row["constraint"]), language)
        c_table.cell(idx, 1).text = f"{float(row['min_gap']):.6f}"
        c_table.cell(idx, 2).text = str(row["worst_model_point"])
        if is_ja:
            c_table.cell(idx, 3).text = "適合" if bool(row["all_ok"]) else "要対応"
        else:
            c_table.cell(idx, 3).text = "PASS" if bool(row["all_ok"]) else "FAIL"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "利源別キャッシュフロー" if is_ja else "Cashflow by Profit Source",
        "全モデルポイント合算の年度分解" if is_ja else "Yearly decomposition across all model points",
    )
    slide.shapes.add_picture(str(cashflow_chart_path), Inches(0.5), Inches(1.5), width=Inches(12.3))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "感応度とリスク" if is_ja else "Sensitivity and Risks",
        "金利・解約率・事業費を±10%変動" if is_ja else "10% shocks on interest, lapse, and expenses",
    )
    s_table = slide.shapes.add_table(
        len(sensitivity_rows) + 1, 6, Inches(0.4), Inches(1.5), Inches(12.2), Inches(3.8)
    ).table
    s_headers = (
        ["シナリオ", "最小IRR", "最小NBV", "最小LSR", "最大PTM", "違反件数"]
        if is_ja
        else ["Scenario", "Min IRR", "Min NBV", "Min LSR", "Max PTM", "Violations"]
    )
    for col, name in enumerate(s_headers):
        s_table.cell(0, col).text = name
    for idx, row in enumerate(sensitivity_rows, start=1):
        s_table.cell(idx, 0).text = _scenario_label(str(row["scenario"]), language)
        s_table.cell(idx, 1).text = _fmt_pct(float(row["min_irr"]), 2)
        s_table.cell(idx, 2).text = f"{float(row['min_nbv']):,.0f}"
        s_table.cell(idx, 3).text = f"{float(row['min_loading_surplus_ratio']):.4f}"
        s_table.cell(idx, 4).text = f"{float(row['max_premium_to_maturity']):.4f}"
        s_table.cell(idx, 5).text = str(int(row["violation_count"]))
    risk_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.0), Inches(1.0))
    risk_tf = risk_box.text_frame
    risk_tf.text = (
        "主なリスク: 高年齢・短期の点でPTMマージンが薄く、事業費上振れでNBVが大きく低下し得ます。"
        if is_ja
        else "Key risk: PTM margin is tight for older/short-term points; "
        "expense up-shock can materially reduce NBV."
    )
    _set_text_style(risk_tf, size=Pt(15), color=RGBColor(40, 40, 40))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "意思決定依頼 / 次アクション" if is_ja else "Decision Ask / Next Actions")
    ask_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.5))
    ask_tf = ask_box.text_frame
    ask_tf.word_wrap = True
    asks = (
        [
            "2枚目の価格テーブルPを承認し、本番見積設定へ反映してください。",
            "現行制約をロックし、毎月または前提変更時に再実行してください。",
            "min_irr < 2.0% または max PTM > 1.056 で即時レビューを起動してください。",
            "ガバナンス審査の統制成果物として out/run_summary_executive.json を利用してください。",
        ]
        if is_ja
        else [
            "Approve the pricing table P in Slide 2 for production quote setup.",
            "Lock current constraints and rerun this report monthly or when assumptions change.",
            "Trigger immediate review if min_irr falls below 2.0% or max PTM exceeds 1.056.",
            "Use out/run_summary_executive.json as the control artifact for governance review.",
        ]
    )
    ask_tf.text = asks[0]
    _set_text_style(ask_tf, size=Pt(20), color=RGBColor(20, 20, 20))
    for item in asks[1:]:
        p = ask_tf.add_paragraph()
        p.text = item
        p.level = 0
        p.runs[0].font.size = Pt(19)
        p.runs[0].font.color.rgb = RGBColor(20, 20, 20)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def report_executive_pptx_from_config(
    config_path: Path,
    *,
    out_path: Path | None = None,
    markdown_path: Path | None = None,
    run_summary_path: Path | None = None,
    deck_out_path: Path | None = None,
    chart_dir: Path | None = None,
    r_start: float = 1.0,
    r_end: float = 1.08,
    r_step: float = 0.005,
    irr_threshold: float = 0.02,
    include_sensitivity: bool = True,
    language: str = "ja",
) -> ExecutiveReportOutputs:
    language = _validate_language(language)
    config_path = config_path.expanduser().resolve()
    config = yaml.safe_load(config_path.read_text(encoding="utf-8"))
    base_dir = resolve_base_dir_from_config(config_path)

    result = run_profit_test(config, base_dir=base_dir)
    execution_context = build_execution_context(
        config=config,
        base_dir=base_dir,
        config_path=config_path,
        command="pricing.cli report-executive-pptx",
        argv=[str(config_path)],
    )
    run_summary = build_run_summary(
        config,
        result,
        source="report_executive_pptx",
        execution_context=execution_context,
    )

    run_summary_output = _resolve_output_path(
        base_dir,
        run_summary_path,
        "out/run_summary_executive.json",
    )
    run_summary_output.parent.mkdir(parents=True, exist_ok=True)
    run_summary_output.write_text(
        json.dumps(run_summary, indent=2, ensure_ascii=True),
        encoding="utf-8",
    )

    deck = build_feasibility_report(
        config=config,
        base_dir=base_dir,
        r_start=float(r_start),
        r_end=float(r_end),
        r_step=float(r_step),
        irr_threshold=float(irr_threshold),
        config_path=config_path,
    )
    deck_output = _resolve_output_path(base_dir, deck_out_path, "out/feasibility_deck_executive.yaml")
    deck_output.parent.mkdir(parents=True, exist_ok=True)
    deck_output.write_text(
        yaml.safe_dump(deck, sort_keys=False),
        encoding="utf-8",
    )

    chart_output_dir = _resolve_output_path(base_dir, chart_dir, "out/charts/executive")
    chart_output_dir.mkdir(parents=True, exist_ok=True)
    agg_cashflow = _aggregate_cashflow(result)
    cashflow_chart = _plot_cashflow_by_profit_source(
        agg_cashflow,
        chart_output_dir / "cashflow_by_profit_source.png",
        language=language,
    )
    premium_chart = _plot_annual_premium_by_model_point(
        result.summary,
        chart_output_dir / "annual_premium_by_model_point.png",
        language=language,
    )

    sensitivity_rows = (
        _build_sensitivity_rows(config, base_dir, chart_output_dir / "sensitivity")
        if include_sensitivity
        else [_scenario_summary("base", config, base_dir)]
    )

    markdown_output = _resolve_output_path(base_dir, markdown_path, "reports/feasibility_report.md")
    markdown_output.parent.mkdir(parents=True, exist_ok=True)
    markdown_output.write_text(
        _build_markdown_report(
            config_path=config_path,
            markdown_path=markdown_output,
            run_summary=run_summary,
            batch_result=result,
            feasibility_deck=deck,
            sensitivity_rows=sensitivity_rows,
            cashflow_chart_path=cashflow_chart,
            premium_chart_path=premium_chart,
            language=language,
        ),
        encoding="utf-8",
    )

    pptx_output = _resolve_output_path(base_dir, out_path, "reports/executive_pricing_deck.pptx")
    constraint_rows = _constraint_status_rows(run_summary)
    _write_executive_pptx(
        out_path=pptx_output,
        run_summary=run_summary,
        summary_df=result.summary.sort_values("model_point"),
        constraint_rows=constraint_rows,
        sensitivity_rows=sensitivity_rows,
        cashflow_chart_path=cashflow_chart,
        premium_chart_path=premium_chart,
        config_path=config_path,
        language=language,
    )

    return ExecutiveReportOutputs(
        pptx_path=pptx_output,
        markdown_path=markdown_output,
        run_summary_path=run_summary_output,
        feasibility_deck_path=deck_output,
        cashflow_chart_path=cashflow_chart,
        premium_chart_path=premium_chart,
    )
